import os
import asyncio
from pptx import Presentation
import json
import openai
import concurrent.futures

async def read_presentation(presentation_path: str):
    presentation = Presentation(presentation_path)
    tasks = []

    for i, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                cleaned_text = shape.text.replace('\t', ' ').strip()
                if cleaned_text:
                    slide_text.append(cleaned_text)

        if slide_text:
            task = asyncio.create_task(process_slide(i, ' '.join(slide_text)))
            tasks.append(task)

    return await asyncio.gather(*tasks)

async def process_slide(slide_number: int, slide_text: str):
    try:
        summary = await slide_summary_by_AI(slide_text)
    except Exception as e:
        summary = f"Error processing slide {slide_number}: {str(e)}"
    return f"Slide {slide_number}:\n{summary}"

def save_to_json(filename: str, data: list):
    with open(filename, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

async def slide_summary_by_AI(slide_text: str):
    api_key = os.getenv('OPENAI_API_KEY')
    if not api_key:
        raise EnvironmentError("OpenAI API key not found. Set it as OPENAI_API_KEY environment variable.")

    messages = [
        {"role": "system", "content": "Summarize for me as briefly and succinctly as possible"},
        {"role": "user", "content": slide_text}
    ]

    loop = asyncio.get_event_loop()
    try:
        with concurrent.futures.ThreadPoolExecutor() as pool:
            completion = await loop.run_in_executor(
                pool,
                lambda: openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=messages)
            )
        response = completion.choices[0].message['content']
    except openai.error.OpenAIError as e:
        response = f"OpenAI API error: {str(e)}"
    except Exception as e:
        response = f"Unexpected error: {str(e)}"

    return response

async def main():
    presentation_path = input("Enter the path to the presentation file: ")
    if os.path.exists(presentation_path):
        file_name, _ = os.path.splitext(os.path.basename(presentation_path))
        summaries = await read_presentation(presentation_path)
        if summaries:
            save_to_json(file_name + ".json", summaries)
            print(f"Summaries saved to {file_name}.json")
        else:
            print(f"No summaries found in {presentation_path}")
    else:
        print(f"The file path does not exist: {presentation_path}")

if __name__ == "__main__":
    asyncio.run(main())